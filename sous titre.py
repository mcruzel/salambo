from pptx import Presentation
import pandas as pd

class config:
    # Define the input PowerPoint file and output Excel file paths
    ppt_path = 'EC4.pptx'  # Replace with the path to your PowerPoint file
    output_excel_path = 'Slide_Notes_Text_Only.xlsx'  # Define output path

# Function to extract notes and format them with separator
def extract_notes_with_separator(ppt_path, output_excel_path):
    # Load the PowerPoint file
    presentation = Presentation(ppt_path)

    # Initialize list to store notes with separators
    notes_text_only = []
    for slide in presentation.slides:
        # Extracting notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note_text = slide.notes_slide.notes_text_frame.text.strip()  # Clean up any extra whitespace
            if note_text:  # If there's any text in notes
                notes_text_only.append({"Note": note_text})  # Append the note
                notes_text_only.append({"Note": "🕓🕓🕓🕓🕓🕓  🕓🕓🕓🕓🕓🕓"})  # Append separator

    # Create DataFrame for the notes
    df_notes_text_only = pd.DataFrame(notes_text_only)

    # Save the DataFrame to an Excel file
    df_notes_text_only.to_excel(output_excel_path, index=False, header=False)

# Run the extraction function
extract_notes_with_separator(config.ppt_path, config.output_excel_path)

print("Extraction completed and saved to", config.output_excel_path)
